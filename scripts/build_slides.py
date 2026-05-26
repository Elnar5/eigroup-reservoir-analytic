"""Build the 26-slide presentation deck for eiGroup submission.

Run from project root:
    python scripts/build_slides_v2.py

Output:
    presentation/eigroup_reservoir_analytics.pptx

This script builds a 16:9 widescreen deck with a warm cream background,
navy + gold accents, and embeds the charts already produced under
outputs/figures/.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu


# -----------------------------------------------------------------------------
# Color palette
# -----------------------------------------------------------------------------
BG_CREAM = RGBColor(0xF5, 0xF1, 0xE8)            # background
NAVY = RGBColor(0x1A, 0x4D, 0x7A)                # primary
GOLD = RGBColor(0xC9, 0xA7, 0x69)                # section accent
INK = RGBColor(0x1A, 0x1A, 0x1A)                 # body dark
MUTED = RGBColor(0x55, 0x55, 0x55)               # body muted
PAGE_NUM = RGBColor(0x99, 0x99, 0x99)            # page number
SUCCESS = RGBColor(0x2D, 0x7D, 0x47)             # zone-good
WARNING = RGBColor(0xB8, 0x86, 0x0B)             # saturation/warning
DANGER = RGBColor(0xC0, 0x39, 0x2B)              # zone-d / negative
ORANGE = RGBColor(0xD6, 0x7D, 0x2C)              # zone-c
ZONE_E = RGBColor(0xA8, 0x39, 0x2B)              # zone-e red
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
HIGHLIGHT_BG = RGBColor(0xFF, 0xF3, 0xD6)        # amber tint for emphasis
TABLE_HEADER_BG = RGBColor(0xF0, 0xEE, 0xE6)
TABLE_ROW_BG = RGBColor(0xFF, 0xFF, 0xFF)
TABLE_BORDER = RGBColor(0xF0, 0xEE, 0xE6)
QUOTE_BG = RGBColor(0xFF, 0xFF, 0xFF)

# -----------------------------------------------------------------------------
# Geometry (16:9 widescreen)
# -----------------------------------------------------------------------------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

MARGIN_L = Inches(0.6)
MARGIN_R = Inches(0.6)
MARGIN_T = Inches(0.45)
MARGIN_B = Inches(0.4)

CONTENT_W = SLIDE_W - MARGIN_L - MARGIN_R    # 12.13"
CONTENT_TOP = Inches(0.95)                    # below section label

# Paths
PROJECT_ROOT = Path(__file__).parent.parent
FIG_DIR = PROJECT_ROOT / "outputs" / "figures"
OUT_PATH = PROJECT_ROOT / "presentation" / "eigroup_reservoir_analytics.pptx"
OUT_PATH.parent.mkdir(parents=True, exist_ok=True)

TOTAL_SLIDES = 26  # updated at end if needed


# -----------------------------------------------------------------------------
# Helper builders
# -----------------------------------------------------------------------------
def add_background(slide, color=BG_CREAM):
    """Fill the entire slide with a solid color rectangle (behind everything)."""
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    # send to back
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return bg


def add_text(
    slide, x, y, w, h, text, *,
    font_size=11, color=INK, bold=False, italic=False,
    align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font_name="Calibri",
    spacing=None,
):
    """Insert a simple text box."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.vertical_anchor = anchor

    p = tf.paragraphs[0]
    p.alignment = align
    if spacing is not None:
        p.line_spacing = spacing
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_name
    return tb


def add_section_label(slide, idx, text, color=GOLD):
    """Top section label like '07 — PART A · THE DISCOVERY'."""
    label = f"{idx:02d}  —  {text}"
    add_text(
        slide, MARGIN_L, MARGIN_T, CONTENT_W, Inches(0.3),
        label, font_size=10, color=color, bold=True,
    )


def add_page_number(slide, idx, total=TOTAL_SLIDES):
    """Bottom-right page number."""
    box = slide.shapes.add_textbox(
        SLIDE_W - Inches(1.0), SLIDE_H - Inches(0.35),
        Inches(0.7), Inches(0.25),
    )
    tf = box.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = f"{idx:02d} / {total:02d}"
    r.font.size = Pt(8)
    r.font.color.rgb = PAGE_NUM
    r.font.name = "Calibri"


def add_title(slide, text, color=INK, font_size=24):
    """Main slide title, just under the section label."""
    add_text(
        slide, MARGIN_L, Inches(0.85), CONTENT_W, Inches(0.6),
        text, font_size=font_size, color=color, bold=True,
    )


def add_subtitle(slide, text, y=Inches(1.5), color=MUTED, font_size=14, italic=False):
    """Subtitle below the title."""
    add_text(
        slide, MARGIN_L, y, CONTENT_W, Inches(0.45),
        text, font_size=font_size, color=color, italic=italic,
    )


def add_big_number(slide, x, y, value, label, *,
                   value_size=54, label_size=10,
                   value_color=NAVY, label_color=MUTED, w=Inches(2.6)):
    """A '3,514 / depth samples' style number block."""
    add_text(slide, x, y, w, Inches(0.85), value,
             font_size=value_size, color=value_color, bold=True, spacing=0.9)
    add_text(slide, x, y + Inches(0.95), w, Inches(0.3), label,
             font_size=label_size, color=label_color)


def add_quote_card(slide, x, y, w, h, *, bg=QUOTE_BG, accent=GOLD, accent_width=Inches(0.04)):
    """A card with a colored left strip — used for highlighted callouts."""
    # Background
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = bg
    box.line.color.rgb = TABLE_BORDER
    box.line.width = Pt(0.5)
    # Accent strip on left
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, accent_width, h)
    strip.fill.solid()
    strip.fill.fore_color.rgb = accent
    strip.line.fill.background()
    return box


def add_image_safely(slide, image_path, x, y, w=None, h=None):
    """Add an image if the file exists, else a placeholder rectangle."""
    p = Path(image_path)
    if p.exists():
        kwargs = {}
        if w is not None:
            kwargs["width"] = w
        if h is not None:
            kwargs["height"] = h
        return slide.shapes.add_picture(str(p), x, y, **kwargs)
    # Placeholder
    placeholder_w = w if w is not None else Inches(6)
    placeholder_h = h if h is not None else Inches(4)
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, placeholder_w, placeholder_h)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
    box.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    tf = box.text_frame
    tf.word_wrap = True
    p_run = tf.paragraphs[0]
    p_run.alignment = PP_ALIGN.CENTER
    p_run.vertical_anchor = MSO_ANCHOR.MIDDLE
    r = p_run.add_run()
    r.text = f"[chart missing: {p.name}]"
    r.font.size = Pt(10)
    r.font.color.rgb = MUTED
    r.font.italic = True
    return box


def add_left_accent_strip(slide, color=NAVY, width=Inches(0.06)):
    """Vertical navy strip on the left side — used on the title slide."""
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, width, SLIDE_H)
    strip.fill.solid()
    strip.fill.fore_color.rgb = color
    strip.line.fill.background()


def add_horizontal_rule(slide, y, *, color=TABLE_BORDER, x=MARGIN_L,
                       w=None, height=Pt(0.5)):
    """Thin horizontal divider."""
    if w is None:
        w = CONTENT_W
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, height)
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()


def style_table_cell(cell, *, text, font_size=10, bold=False, color=INK,
                     bg=None, align=PP_ALIGN.LEFT):
    """Style a python-pptx table cell."""
    cell.text = ""  # clear
    tf = cell.text_frame
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(font_size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.name = "Calibri"
    if bg is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg


# -----------------------------------------------------------------------------
# Per-slide builders
# -----------------------------------------------------------------------------
def build_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    blank = prs.slide_layouts[6]  # Blank

    def new_slide(idx, section_text=None):
        slide = prs.slides.add_slide(blank)
        add_background(slide)
        if section_text:
            add_section_label(slide, idx, section_text)
        add_page_number(slide, idx)
        return slide

    # === SLIDE 1 — Title ===
    s = prs.slides.add_slide(blank)
    add_background(s)
    add_left_accent_strip(s, NAVY)
    add_text(s, Inches(0.8), Inches(2.4), Inches(8),  Inches(0.35),
             "RESERVOIR ANALYTICS  ·  MAY 2026", font_size=11,
             color=GOLD, bold=True)
    add_text(s, Inches(0.8), Inches(2.8), Inches(11), Inches(1.8),
             "Beyond the\nHeadline Numbers", font_size=54,
             color=INK, bold=True, spacing=1.0)
    add_text(s, Inches(0.82), Inches(5.0), Inches(11), Inches(0.5),
             "A petrophysical investigation of 18,167 depth samples.",
             font_size=18, color=MUTED, italic=True)
    add_horizontal_rule(s, Inches(5.9), x=Inches(0.82), w=Inches(2.5),
                        color=GOLD, height=Pt(1.5))
    add_text(s, Inches(0.82), Inches(6.1), Inches(11), Inches(0.35),
             "Kamil Muradli", font_size=14, color=INK, bold=True)
    add_text(s, Inches(0.82), Inches(6.45), Inches(11), Inches(0.3),
             "eiGroup Associate Data Scientist Assessment",
             font_size=11, color=MUTED)
    add_text(s, Inches(0.82), Inches(6.75), Inches(11), Inches(0.3),
             "Submitted May 2026", font_size=10, color=PAGE_NUM)

    # === SLIDE 2 — The Assignment ===
    s = new_slide(2, "OVERVIEW  ·  THE ASSIGNMENT")
    add_title(s, "Four parts. One coherent investigation.")
    add_subtitle(s, "Each part both stands alone and informs the next.")

    parts = [
        ("Part A", "Data Quality",
         "Load, join, validate. Discover what the dataset hides."),
        ("Part B", "Per-Zone Metrics",
         "Five required + seven bonus diagnostics across 35 (well, zone) groups."),
        ("Part C", "Cutoff Sensitivity & Field Views",
         "Sweep cutoffs to bound conclusions. Visualize the field from six angles."),
        ("Part D", "Sub-Zone Definition",
         "Find geologically meaningful sub-zones consistent across wells."),
    ]
    col_w = Inches(2.95)
    gap = Inches(0.15)
    start_x = MARGIN_L
    y0 = Inches(2.3)
    for i, (tag, name, desc) in enumerate(parts):
        x = start_x + (col_w + gap) * i
        # Card
        card = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y0, col_w, Inches(3.5))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = TABLE_BORDER
        card.line.width = Pt(0.5)
        # Gold strip top
        strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y0, col_w, Inches(0.06))
        strip.fill.solid()
        strip.fill.fore_color.rgb = GOLD
        strip.line.fill.background()
        # Tag
        add_text(s, x + Inches(0.25), y0 + Inches(0.3), col_w - Inches(0.5), Inches(0.3),
                 tag, font_size=10, color=GOLD, bold=True)
        # Name
        add_text(s, x + Inches(0.25), y0 + Inches(0.65), col_w - Inches(0.5), Inches(0.6),
                 name, font_size=18, color=NAVY, bold=True, spacing=1.0)
        # Desc
        add_text(s, x + Inches(0.25), y0 + Inches(1.45), col_w - Inches(0.5), Inches(1.8),
                 desc, font_size=11, color=MUTED, spacing=1.3)

    add_text(s, MARGIN_L, Inches(6.5), CONTENT_W, Inches(0.4),
             "→  Five required deliverables, four exceeded. Every part backed by tests and reproducible CLI.",
             font_size=11, color=NAVY, italic=True)

    # === SLIDE 3 — Our Approach ===
    s = new_slide(3, "OVERVIEW  ·  PIPELINE ARCHITECTURE")
    add_title(s, "One command per part. Reproducible end-to-end.")
    add_subtitle(s, "Each CLI command writes parquet caches, CSV deliverables, and Markdown reports.")
    # Embed the architecture diagram — sized to leave room for tagline below
    arch_path = FIG_DIR / "00_architecture.png"
    add_image_safely(s, arch_path, Inches(1.0), Inches(2.0), w=Inches(11.3))

    # === SLIDE 4 — The Dataset ===
    s = new_slide(4, "OVERVIEW  ·  THE DATASET")
    add_title(s, "Seven wells. Five zones. Eighteen thousand samples.")
    add_subtitle(s, "Five petrophysical logs per sample: vsh, phit, perm, sw, depth.")

    # 4 metric cards
    y = Inches(2.4)
    card_w = Inches(2.85)
    gap = Inches(0.15)
    metrics = [
        ("18,167", "depth samples"),
        ("7", "wells"),
        ("5", "zones (A–E)"),
        ("0.2 – 0.5 m", "sampling step"),
    ]
    for i, (val, lbl) in enumerate(metrics):
        x = MARGIN_L + (card_w + gap) * i
        card = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_w, Inches(2.0))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = TABLE_BORDER
        card.line.width = Pt(0.5)
        add_text(s, x + Inches(0.3), y + Inches(0.4), card_w - Inches(0.6),
                 Inches(0.9), val, font_size=36, color=NAVY, bold=True, spacing=1.0)
        add_text(s, x + Inches(0.3), y + Inches(1.35), card_w - Inches(0.6),
                 Inches(0.3), lbl, font_size=11, color=MUTED)

    # Bottom info
    add_horizontal_rule(s, Inches(5.0))
    add_text(s, MARGIN_L, Inches(5.15), CONTENT_W, Inches(0.35),
             "Quality issues by design",
             font_size=12, color=NAVY, bold=True)
    add_text(s, MARGIN_L, Inches(5.55), CONTENT_W, Inches(0.4),
             "•   Well 5 sampled at 0.5 m vs 0.2 m elsewhere   ·   "
             "Well 3 has 78 NaN porosity values   ·   "
             "3,514 samples sit at the perm tool ceiling",
             font_size=11, color=MUTED, spacing=1.4)
    add_text(s, MARGIN_L, Inches(6.6), CONTENT_W, Inches(0.4),
             "→  All three were uncovered, documented, and handled. Detail in the next four slides.",
             font_size=11, color=NAVY, italic=True)

    # === SLIDE 5 — Part A: Data Integration ===
    s = new_slide(5, "PART A  ·  DATA INTEGRATION")
    add_title(s, "Joining 18,167 samples to 35 zone tops.")
    add_subtitle(s, "Pandas merge_asof on per-well depth ranges, with boundary validation.")

    # Left: pipeline steps
    y0 = Inches(2.3)
    steps = [
        ("01", "Load 7 well CSVs + zones lookup",
         "Schema validation. well_id column added per file."),
        ("02", "merge_asof on depth (per well)",
         "direction='backward', by='well_id'. Each sample assigned to its containing zone."),
        ("03", "Compute per-well dz",
         "depth.diff() per well_id. Critical: well_5 gets dz=0.5, others dz=0.2."),
        ("04", "Validate boundaries + run QC",
         "Range checks, NaN counts, saturation flag. Three real issues found."),
    ]
    sx = MARGIN_L
    for i, (num, head, body) in enumerate(steps):
        sy = y0 + Inches(1.05) * i
        # Step number
        add_text(s, sx, sy, Inches(0.5), Inches(0.6),
                 num, font_size=22, color=GOLD, bold=True)
        # Headline
        add_text(s, sx + Inches(0.65), sy, Inches(11), Inches(0.4),
                 head, font_size=14, color=NAVY, bold=True)
        # Detail
        add_text(s, sx + Inches(0.65), sy + Inches(0.4), Inches(11), Inches(0.5),
                 body, font_size=10, color=MUTED, spacing=1.3)

    add_text(s, MARGIN_L, Inches(6.8), CONTENT_W, Inches(0.3),
             "→  Output: master_table.parquet (18,167 rows × 9 columns). Single source of truth.",
             font_size=11, color=NAVY, italic=True)

    # === SLIDE 6 — Three Quality Issues ===
    s = new_slide(6, "PART A  ·  WHAT WE FOUND")
    add_title(s, "Three quality issues by design. All flagged.")
    add_subtitle(s, "The dataset was deliberately constructed with these. We discovered all three.")

    y0 = Inches(2.4)
    card_h = Inches(3.6)
    card_w = Inches(3.95)
    gap = Inches(0.15)
    issues = [
        ("Tool saturation",
         "3,514", "samples",
         "Permeability values pinned at exactly 15,000 mD — the upper measurement limit.",
         "19.34% of all samples",
         WARNING),
        ("Well 5 sampling",
         "0.5 m", "step",
         "Six wells sampled at 0.2 m; well 5 at 0.5 m. 2.5× difference per sample.",
         "Solved by depth-weighting everywhere",
         NAVY),
        ("Well 3 missing data",
         "78", "NaN phit",
         "Distributed across zones in well 3. Consistent with tool malfunction at depths.",
         "0.43% of total samples",
         GOLD),
    ]
    for i, (head, big, lbl, body, footer, color) in enumerate(issues):
        x = MARGIN_L + (card_w + gap) * i
        card = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y0, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = TABLE_BORDER
        card.line.width = Pt(0.5)
        # accent strip top
        strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y0, card_w, Inches(0.08))
        strip.fill.solid()
        strip.fill.fore_color.rgb = color
        strip.line.fill.background()
        # Title
        add_text(s, x + Inches(0.3), y0 + Inches(0.3), card_w - Inches(0.6),
                 Inches(0.4), head, font_size=12, color=MUTED, bold=True)
        # Big number
        add_text(s, x + Inches(0.3), y0 + Inches(0.75), card_w - Inches(0.6),
                 Inches(0.9), big, font_size=40, color=color, bold=True, spacing=1.0)
        # Label
        add_text(s, x + Inches(0.3), y0 + Inches(1.65), card_w - Inches(0.6),
                 Inches(0.3), lbl, font_size=10, color=MUTED)
        # Body
        add_text(s, x + Inches(0.3), y0 + Inches(2.05), card_w - Inches(0.6),
                 Inches(1.1), body, font_size=10, color=INK, spacing=1.4)
        # Footer
        add_text(s, x + Inches(0.3), y0 + Inches(3.15), card_w - Inches(0.6),
                 Inches(0.3), footer, font_size=10, color=color, italic=True)

    # === SLIDE 7 — The Discovery (Hero slide) ===
    s = new_slide(7, "PART A  ·  THE DISCOVERY")
    add_title(s, "A number that shouldn't exist.")

    # Big number centered
    add_text(s, MARGIN_L, Inches(1.8), CONTENT_W, Inches(2.3),
             "3,514", font_size=160, color=NAVY, bold=True,
             align=PP_ALIGN.CENTER, spacing=1.0)
    add_text(s, MARGIN_L, Inches(4.4), CONTENT_W, Inches(0.45),
             "of 18,167 samples report exactly 15,000 mD",
             font_size=18, color=INK, align=PP_ALIGN.CENTER, italic=True)
    add_text(s, MARGIN_L, Inches(4.95), CONTENT_W, Inches(0.4),
             "— to machine precision. Not 14,985. Not 15,012. Exactly 15,000.",
             font_size=13, color=MUTED, align=PP_ALIGN.CENTER)

    # Bottom callout
    add_quote_card(s, Inches(2.5), Inches(5.85), Inches(8.35), Inches(0.8),
                   bg=HIGHLIGHT_BG, accent=GOLD)
    add_text(s, Inches(2.75), Inches(5.95), Inches(8), Inches(0.6),
             "A real measurement near a tool's ceiling is continuous. A discrete spike at exactly the maximum is the signature of saturation.",
             font_size=11, color=INK, italic=True, spacing=1.4)

    # === SLIDE 8 — What It Changes ===
    s = new_slide(8, "PART A  ·  CONSEQUENCES OF THE DISCOVERY")
    add_title(s, "Four downstream effects we had to handle.")
    add_subtitle(s, "Every part of the pipeline below treats saturation as a first-class concern.")

    y0 = Inches(2.4)
    items = [
        ("kh estimates",
         "Lower bounds, not best estimates.",
         "Saturated samples contribute 15,000 × dz to kh. Real perm is higher. Effect surfaces in every kh number reported."),
        ("Average permeability",
         "Dragged down toward the ceiling.",
         "Arithmetic mean of 14,997 mD in Zone B is essentially the ceiling — not the true mean."),
        ("Lorenz heterogeneity",
         "Collapses to ~0 when saturation is severe.",
         "Zone B's Lorenz ≈ 0.001 is mathematical inevitability, not a homogeneity finding. The instrument can't see the heterogeneity."),
        ("Clustering on perm",
         "Fails silently on saturated features.",
         "Part D Zone B: two of three sub-zones get log_perm centroids of 4.175938 vs 4.176091 — indistinguishable."),
    ]
    row_h = Inches(1.0)
    for i, (head, summary, body) in enumerate(items):
        ry = y0 + row_h * i
        # number / dot
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, MARGIN_L, ry + Inches(0.18),
                                  Inches(0.16), Inches(0.16))
        dot.fill.solid()
        dot.fill.fore_color.rgb = GOLD
        dot.line.fill.background()
        # head
        add_text(s, MARGIN_L + Inches(0.4), ry, Inches(3.0), Inches(0.4),
                 head, font_size=14, color=NAVY, bold=True)
        # summary
        add_text(s, MARGIN_L + Inches(3.5), ry, Inches(8.0), Inches(0.4),
                 summary, font_size=12, color=INK, bold=True)
        # body
        add_text(s, MARGIN_L + Inches(3.5), ry + Inches(0.4), Inches(8.6),
                 Inches(0.5), body, font_size=10, color=MUTED, spacing=1.3)

    add_text(s, MARGIN_L, Inches(6.8), CONTENT_W, Inches(0.3),
             "→  We did not drop saturated samples. We kept, counted, and surfaced them everywhere.",
             font_size=11, color=NAVY, italic=True)

    # === SLIDE 9 — Part B intro ===
    s = new_slide(9, "PART B  ·  PER-ZONE METRICS")
    add_title(s, "Twelve metrics per (well, zone). Thirty-five rows.")
    add_subtitle(s, "Five required by the case + seven bonus diagnostics we added.")

    # Left column: required
    add_text(s, MARGIN_L, Inches(2.3), Inches(5.9), Inches(0.4),
             "REQUIRED  ·  case statement",
             font_size=10, color=GOLD, bold=True)
    required = [
        ("Gross thickness", "sum(dz) over zone"),
        ("Net thickness", "sum(dz) where vsh ≤ 0.5 AND phit ≥ 0.08"),
        ("Avg phit (net)", "mean(phit) over net rows"),
        ("Avg perm (net)", "mean(perm) over net rows"),
        ("kh (flow capacity)", "sum(perm × dz) over net rows"),
    ]
    rx = MARGIN_L
    ry = Inches(2.75)
    for name, formula in required:
        add_text(s, rx, ry, Inches(2.4), Inches(0.3), name,
                 font_size=11, color=NAVY, bold=True)
        add_text(s, rx + Inches(2.4), ry, Inches(3.5), Inches(0.3),
                 formula, font_size=10, color=MUTED)
        ry += Inches(0.5)

    # Right column: bonus
    add_text(s, Inches(6.8), Inches(2.3), Inches(6.0), Inches(0.4),
             "BONUS  ·  earns its place",
             font_size=10, color=GOLD, bold=True)
    bonus = [
        ("NTG",            "net/gross. Standard quality KPI."),
        ("kh-weighted perm","Engineering-correct mean."),
        ("Lorenz coefficient","Heterogeneity in [0,1]. Motivates Part D."),
        ("n_samples_net",   "Sanity check."),
        ("n_phit_nan",      "NaN exclusions counted (well 3)."),
        ("n_perm_saturated_in_net", "Tool-cap count. Makes lower-bound visible."),
        ("frac_saturated",  "Saturation as a fraction of net."),
    ]
    bx = Inches(6.8)
    by = Inches(2.75)
    for name, why in bonus:
        add_text(s, bx, by, Inches(2.6), Inches(0.3), name,
                 font_size=11, color=NAVY, bold=True)
        add_text(s, bx + Inches(2.6), by, Inches(3.5), Inches(0.3),
                 why, font_size=10, color=MUTED)
        by += Inches(0.4)

    # === SLIDE 10 — Heatmap ===
    s = new_slide(10, "PART B  ·  WHERE THE FLOW IS")
    add_title(s, "kh by well × zone. The whole field on one grid.")
    add_subtitle(s, "Log-coloured cells. Saturated-sample counts annotated where present.")
    add_image_safely(s, FIG_DIR / "01_kh_heatmap.png",
                     Inches(1.5), Inches(2.1), w=Inches(10.3))
    add_text(s, MARGIN_L, Inches(6.85), CONTENT_W, Inches(0.35),
             "→  Zone B uniformly dominant across all wells. Zone D uniformly weak. Well 7 leads at 2.37 M mD·m.",
             font_size=11, color=NAVY, italic=True)

    # === SLIDE 11 — 5 zone signatures ===
    s = new_slide(11, "PART B  ·  FIVE ZONE SIGNATURES")
    add_title(s, "Five zones. Five distinct stories.")
    add_subtitle(s, "Field-level rollup of net thickness, kh, and saturation.")

    # Build a table
    rows, cols = 6, 7
    tbl_x = MARGIN_L
    tbl_y = Inches(2.4)
    tbl_w = CONTENT_W
    tbl_h = Inches(3.3)
    tbl = s.shapes.add_table(rows, cols, tbl_x, tbl_y, tbl_w, tbl_h).table
    # column widths
    widths = [0.6, 1.05, 1.05, 1.3, 1.05, 1.4, 5.65]
    total = sum(widths)
    for i, w in enumerate(widths):
        tbl.columns[i].width = Inches(w * 12.13 / total)

    headers = ["Zone", "Net (m)", "NTG", "kh (mD·m)", "Avg perm", "Saturated", "Signature"]
    for c, h in enumerate(headers):
        style_table_cell(tbl.cell(0, c), text=h, font_size=10, bold=True,
                         color=NAVY, bg=TABLE_HEADER_BG)

    zone_data = [
        ("A", "466.0", "0.63", "267 K", "572 mD", "0",
         "Clean top reservoir — most cutoff-sensitive", SUCCESS),
        ("B", "711.9", "0.93", "10.7 M  ⚠", "14,997 mD", "3,328",
         "Dominant flow zone — every kh is a lower bound", NAVY),
        ("C", "915.0", "0.84", "680 K", "743 mD", "17",
         "Heterogeneous (Lorenz 0.65) — motivates Part D clustering", ORANGE),
        ("D", "52.9", "0.10", "42", "0.79 mD", "0",
         "Tight rock — non-reservoir at any cutoff", DANGER),
        ("E", "586.8", "0.70", "1.2 M", "2,045 mD", "39",
         "Deep reservoir — most defensible high-perm zone", ZONE_E),
    ]
    for r, (z, net, ntg, kh, perm, sat, sig, color) in enumerate(zone_data, start=1):
        # alternating row background subtle
        bg = TABLE_ROW_BG if r % 2 == 1 else RGBColor(0xFA, 0xF8, 0xF2)
        style_table_cell(tbl.cell(r, 0), text=z, font_size=12, bold=True,
                         color=color, bg=bg, align=PP_ALIGN.CENTER)
        style_table_cell(tbl.cell(r, 1), text=net, font_size=11, color=INK, bg=bg, align=PP_ALIGN.RIGHT)
        style_table_cell(tbl.cell(r, 2), text=ntg, font_size=11, color=INK, bg=bg, align=PP_ALIGN.RIGHT)
        kh_color = WARNING if "⚠" in kh else INK
        style_table_cell(tbl.cell(r, 3), text=kh, font_size=11, color=kh_color, bg=bg,
                         bold="⚠" in kh, align=PP_ALIGN.RIGHT)
        style_table_cell(tbl.cell(r, 4), text=perm, font_size=11, color=INK, bg=bg, align=PP_ALIGN.RIGHT)
        sat_color = WARNING if int(sat.replace(",", "")) > 100 else INK
        style_table_cell(tbl.cell(r, 5), text=sat, font_size=11, color=sat_color,
                         bg=bg, bold=int(sat.replace(",", "")) > 100, align=PP_ALIGN.RIGHT)
        style_table_cell(tbl.cell(r, 6), text=sig, font_size=10, color=MUTED, bg=bg)

    add_text(s, MARGIN_L, Inches(6.4), CONTENT_W, Inches(0.4),
             "→  Zone B's headline kh is a lower bound; Zone E is the largest defensible kh in the field at 1.2 M.",
             font_size=11, color=NAVY, italic=True)

    # === SLIDE 12 — Well 7 Paradox ===
    s = new_slide(12, "PART B  ·  THE WELL 7 PARADOX")
    add_title(s, "Well 7 leads the kh ranking. But also leads in censoring.")
    add_subtitle(s, "The two go together — and the implication is uncomfortable.")

    y0 = Inches(2.4)
    # Two big numbers
    add_big_number(s, Inches(1.0), y0, "2.59 M",
                   "mD·m total kh (rank 1 of 7)",
                   value_size=60, value_color=NAVY)
    add_big_number(s, Inches(5.7), y0, "798",
                   "saturated samples (30% of well 7)",
                   value_size=60, value_color=WARNING)
    add_big_number(s, Inches(9.7), y0, "?",
                   "true gap to the rest of the field",
                   value_size=60, value_color=DANGER, w=Inches(3))

    add_horizontal_rule(s, Inches(5.0))

    add_text(s, MARGIN_L, Inches(5.2), CONTENT_W, Inches(0.5),
             "Saturation makes apparent kh larger. The bound is loosest where the headline number is biggest.",
             font_size=14, color=INK, italic=True, spacing=1.3)
    add_quote_card(s, MARGIN_L, Inches(6.0), CONTENT_W, Inches(0.9), accent=NAVY)
    add_text(s, MARGIN_L + Inches(0.2), Inches(6.1), CONTENT_W - Inches(0.4), Inches(0.8),
             "Recommendation: any well-to-well kh comparison must report saturation fraction alongside the kh number, otherwise rankings are partly an instrument artefact.",
             font_size=11, color=INK, spacing=1.4)

    # === SLIDE 13 — Part C.1 intro ===
    s = new_slide(13, "PART C.1  ·  CUTOFF SENSITIVITY")
    add_title(s, "Nine cutoffs. 315 results. One robust picture.")
    add_subtitle(s, "Default vsh = 0.5 is a literature value, not calibrated. Sweep instead.")

    y0 = Inches(2.4)
    # Why sweep
    add_quote_card(s, MARGIN_L, y0, CONTENT_W, Inches(0.8), accent=GOLD)
    add_text(s, MARGIN_L + Inches(0.2), y0 + Inches(0.1), CONTENT_W - Inches(0.4),
             Inches(0.6),
             "Without core or production data to calibrate against, picking one cutoff is a guess. The honest move: bound the conclusions by sweeping the threshold.",
             font_size=12, color=INK, spacing=1.4, italic=True)
    # Sweep grid
    y1 = Inches(3.5)
    add_text(s, MARGIN_L, y1, CONTENT_W, Inches(0.3),
             "Sweep parameters", font_size=10, color=GOLD, bold=True)
    grid_top = y1 + Inches(0.4)
    grid_items = [
        ("9", "vsh cutoffs"),
        ("0.30 → 0.70", "range (step 0.05)"),
        ("315", "result rows"),
        ("600", "bootstrap resamples"),
    ]
    for i, (val, lbl) in enumerate(grid_items):
        x = MARGIN_L + Inches(3.05) * i
        add_text(s, x, grid_top, Inches(3.0), Inches(0.85),
                 val, font_size=36, color=NAVY, bold=True, spacing=1.0)
        add_text(s, x, grid_top + Inches(0.95), Inches(3.0), Inches(0.3),
                 lbl, font_size=10, color=MUTED)

    add_text(s, MARGIN_L, Inches(6.7), CONTENT_W, Inches(0.4),
             "→  Next slide: the picture that bounds every conclusion in this submission.",
             font_size=11, color=NAVY, italic=True)

    # === SLIDE 14 — Sensitivity curves ===
    s = new_slide(14, "PART C.1  ·  SENSITIVITY CURVES")
    add_title(s, "Two flat zones bracket three sensitive ones.")
    add_subtitle(s, "NTG_field vs vsh cutoff. Bold = zone average. Thin = per-well trajectory.")
    add_image_safely(s, FIG_DIR / "04_ntg_sensitivity.png",
                     Inches(2.1), Inches(2.1), w=Inches(9.1))
    add_text(s, MARGIN_L, Inches(6.85), CONTENT_W, Inches(0.35),
             "→  Zone B robust at top (79-93%). Zone D robust failure (≤29%). A, C, E swing 50-79 pp.",
             font_size=11, color=NAVY, italic=True)

    # === SLIDE 15 — Robust vs cutoff-driven ===
    s = new_slide(15, "PART C.1  ·  ROBUST vs CUTOFF-DRIVEN")
    add_title(s, "Which findings survive any reasonable cutoff?")
    add_subtitle(s, "Five findings hold robustly. Three depend on the threshold.")

    y0 = Inches(2.4)
    # Two columns
    col_w = Inches(5.9)
    col_gap = Inches(0.35)
    # Left — Robust
    left_x = MARGIN_L
    head_box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_x, y0, col_w, Inches(0.5))
    head_box.fill.solid()
    head_box.fill.fore_color.rgb = SUCCESS
    head_box.line.fill.background()
    add_text(s, left_x + Inches(0.2), y0 + Inches(0.1), col_w - Inches(0.4),
             Inches(0.3), "ROBUST  ·  survive every cutoff",
             font_size=12, color=WHITE, bold=True)
    robust = [
        "Zone B is the field's dominant flow zone",
        "Zone B kh is a lower bound (saturation)",
        "Zone D is effectively non-reservoir",
        "Well 7 leads on visible kh",
        "Zone E is the most defensible high-perm zone",
    ]
    for i, item in enumerate(robust):
        ry = y0 + Inches(0.65) + Inches(0.55) * i
        # check mark
        add_text(s, left_x + Inches(0.2), ry, Inches(0.3), Inches(0.4),
                 "✓", font_size=14, color=SUCCESS, bold=True)
        add_text(s, left_x + Inches(0.55), ry, col_w - Inches(0.7),
                 Inches(0.4), item, font_size=11, color=INK, spacing=1.3)

    # Right — Cutoff-driven
    right_x = left_x + col_w + col_gap
    head_box2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_x, y0, col_w, Inches(0.5))
    head_box2.fill.solid()
    head_box2.fill.fore_color.rgb = WARNING
    head_box2.line.fill.background()
    add_text(s, right_x + Inches(0.2), y0 + Inches(0.1), col_w - Inches(0.4),
             Inches(0.3), "CUTOFF-DRIVEN  ·  depend on threshold",
             font_size=12, color=WHITE, bold=True)
    cutoff_driven = [
        ("Zone A's NTG = 0.63", "ranges 0.15 to 0.94"),
        ("Zone C's NTG = 0.84", "ranges 0.38 to 0.94"),
        ("Zone E's NTG = 0.70", "ranges 0.32 to 0.89"),
    ]
    for i, (head, det) in enumerate(cutoff_driven):
        ry = y0 + Inches(0.65) + Inches(0.55) * i
        add_text(s, right_x + Inches(0.2), ry, Inches(0.3), Inches(0.4),
                 "△", font_size=14, color=WARNING, bold=True)
        add_text(s, right_x + Inches(0.55), ry, col_w - Inches(0.7),
                 Inches(0.4),
                 head, font_size=11, color=INK, bold=True, spacing=1.3)
        add_text(s, right_x + Inches(0.55), ry + Inches(0.25),
                 col_w - Inches(0.7), Inches(0.4),
                 det, font_size=10, color=MUTED)

    add_text(s, MARGIN_L, Inches(6.7), CONTENT_W, Inches(0.4),
             "→  For cutoff-driven findings, we report ranges. For robust findings, point estimates are defensible.",
             font_size=11, color=NAVY, italic=True)

    # === SLIDE 16 — Lorenz curves ===
    s = new_slide(16, "PART C.2  ·  LORENZ HETEROGENEITY")
    add_title(s, "The signal that motivated Part D.")
    add_subtitle(s, "Flow concentration per zone. 0 = uniform, 1 = single super-streak.")
    add_image_safely(s, FIG_DIR / "05_lorenz_curves.png",
                     Inches(0.5), Inches(2.1), w=Inches(7.5))
    # Right column with key values
    rx = Inches(8.4)
    ry = Inches(2.3)
    add_text(s, rx, ry, Inches(4.5), Inches(0.3),
             "Lorenz coefficients", font_size=10, color=GOLD, bold=True)
    lorenz_vals = [
        ("Zone C", "0.65", "real heterogeneity — needs sub-zoning", ORANGE),
        ("Zone E", "~0.52", "moderate, typical clean sand", ZONE_E),
        ("Zone A", "~0.45", "moderate", SUCCESS),
        ("Zone D", "0.30-0.48", "noisy, small samples", DANGER),
        ("Zone B", "~0.001", "saturation artefact", NAVY),
    ]
    ry_start = Inches(2.7)
    for i, (zone, val, note, color) in enumerate(lorenz_vals):
        cy = ry_start + Inches(0.75) * i
        add_text(s, rx, cy, Inches(1.0), Inches(0.3),
                 zone, font_size=12, color=color, bold=True)
        add_text(s, rx + Inches(1.05), cy, Inches(1.2), Inches(0.3),
                 val, font_size=14, color=NAVY, bold=True)
        add_text(s, rx, cy + Inches(0.35), Inches(4.5), Inches(0.3),
                 note, font_size=10, color=MUTED, italic=True)

    add_text(s, MARGIN_L, Inches(6.85), CONTENT_W, Inches(0.35),
             "→  Zone C's 0.65 is the highest non-artefactual value. That's the bridge to Part D.",
             font_size=11, color=NAVY, italic=True)

    # === SLIDE 17 — Box plot ===
    s = new_slide(17, "PART C.2  ·  ANSWERING THE CASE PROMPT")
    add_title(s, "Which zones are consistently strong or weak?")
    add_subtitle(s, "NTG and log(kh) distributions across 7 wells, per zone.")
    add_image_safely(s, FIG_DIR / "09_zone_quality_boxplot.png",
                     Inches(1.5), Inches(2.1), w=Inches(10.3))
    add_text(s, MARGIN_L, Inches(6.85), CONTENT_W, Inches(0.35),
             "→  Zone B tight at top (consistent strong). Zone D tight at bottom (consistent weak). Zone A widest box (variable).",
             font_size=11, color=NAVY, italic=True)

    # === SLIDE 18 — Part D intro ===
    s = new_slide(18, "PART D  ·  SUB-ZONE DEFINITION")
    add_title(s, "Pick one zone. Subdivide it. Consistent across wells.")
    add_subtitle(s, "We picked Zone C. We also ran Zone B as a controlled negative result.")

    # Two columns - why C, method
    y0 = Inches(2.4)
    # Left — Why Zone C
    add_text(s, MARGIN_L, y0, Inches(5.9), Inches(0.4),
             "WHY ZONE C", font_size=11, color=GOLD, bold=True)
    why_c = [
        ("Lorenz 0.65", "highest non-artefactual heterogeneity"),
        ("1,094 m gross", "thickest zone in the field"),
        ("Only 0.4% saturated", "instrument actually sees the perm range"),
    ]
    for i, (head, det) in enumerate(why_c):
        cy = y0 + Inches(0.5) + Inches(0.85) * i
        add_text(s, MARGIN_L, cy, Inches(5.9), Inches(0.3),
                 head, font_size=14, color=NAVY, bold=True)
        add_text(s, MARGIN_L, cy + Inches(0.35), Inches(5.9), Inches(0.3),
                 det, font_size=11, color=MUTED, italic=True)
    # Right — Method
    rx = Inches(7.0)
    add_text(s, rx, y0, Inches(5.9), Inches(0.4),
             "METHOD", font_size=11, color=GOLD, bold=True)
    method = [
        ("K-Means on 6 features", "vsh, phit, log(perm), sw, eff_φ, hc_φ"),
        ("Standardized + pooled fit", "z-score then fit on all 7 wells together"),
        ("GMM cross-check", "second algorithm for structure validation"),
    ]
    for i, (head, det) in enumerate(method):
        cy = y0 + Inches(0.5) + Inches(0.85) * i
        add_text(s, rx, cy, Inches(5.9), Inches(0.3),
                 head, font_size=14, color=NAVY, bold=True)
        add_text(s, rx, cy + Inches(0.35), Inches(5.9), Inches(0.3),
                 det, font_size=11, color=MUTED, italic=True)

    add_text(s, MARGIN_L, Inches(6.7), CONTENT_W, Inches(0.4),
             "→  The pooled fit is what makes 'same sub-zone across wells' operationally meaningful.",
             font_size=11, color=NAVY, italic=True)

    # === SLIDE 19 — k=3 vs k=2 ===
    s = new_slide(19, "PART D  ·  AN HONEST MODELLING CHOICE")
    add_title(s, "k=2 wins by silhouette. We chose k=3 anyway.")
    add_subtitle(s, "The case asks for 2–3 sub-zones. We document why three is the geological answer.")

    y0 = Inches(2.4)
    # Two boxes
    box_w = Inches(5.9)
    # k=2 box
    box1 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN_L, y0, box_w, Inches(2.0))
    box1.fill.solid()
    box1.fill.fore_color.rgb = WHITE
    box1.line.color.rgb = TABLE_BORDER
    box1.line.width = Pt(0.5)
    add_text(s, MARGIN_L + Inches(0.3), y0 + Inches(0.2), box_w - Inches(0.6),
             Inches(0.4), "k=2  ·  silhouette winner",
             font_size=11, color=MUTED, bold=True)
    add_text(s, MARGIN_L + Inches(0.3), y0 + Inches(0.6), box_w - Inches(0.6),
             Inches(0.8), "0.394", font_size=48, color=MUTED, bold=True, spacing=1.0)
    add_text(s, MARGIN_L + Inches(0.3), y0 + Inches(1.5), box_w - Inches(0.6),
             Inches(0.5),
             "Highest silhouette score, but a coarse 'good vs poor' split. Two operational tiers.",
             font_size=10, color=MUTED, italic=True, spacing=1.3)

    # k=3 box
    box2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.0), y0, box_w, Inches(2.0))
    box2.fill.solid()
    box2.fill.fore_color.rgb = WHITE
    box2.line.color.rgb = NAVY
    box2.line.width = Pt(1.5)
    add_text(s, Inches(7.3), y0 + Inches(0.2), box_w - Inches(0.6), Inches(0.4),
             "k=3  ·  geological winner",
             font_size=11, color=NAVY, bold=True)
    add_text(s, Inches(7.3), y0 + Inches(0.6), box_w - Inches(0.6), Inches(0.8),
             "0.276", font_size=48, color=NAVY, bold=True, spacing=1.0)
    add_text(s, Inches(7.3), y0 + Inches(1.5), box_w - Inches(0.6), Inches(0.5),
             "Three coherent tiers (poor / moderate / best). Three operational decisions.",
             font_size=10, color=INK, italic=True, spacing=1.3)

    # The reason
    add_horizontal_rule(s, Inches(4.85))
    add_text(s, MARGIN_L, Inches(5.0), CONTENT_W, Inches(0.4),
             "Why three: every feature orders monotonically",
             font_size=12, color=GOLD, bold=True)
    add_text(s, MARGIN_L, Inches(5.45), CONTENT_W, Inches(0.4),
             "vsh decreases (0.51 → 0.33 → 0.22)  ·  phit increases (0.14 → 0.20 → 0.25)  ·  perm increases (71 → 231 → 744 mD)",
             font_size=11, color=MUTED, spacing=1.4)
    add_text(s, MARGIN_L, Inches(5.85), CONTENT_W, Inches(0.4),
             "sw decreases (0.62 → 0.47 → 0.32) — hydrocarbon content improves",
             font_size=11, color=MUTED)

    add_text(s, MARGIN_L, Inches(6.8), CONTENT_W, Inches(0.4),
             "→  Random or spurious clusters do not order monotonically on multiple petrophysical features.",
             font_size=11, color=NAVY, italic=True)

    # === SLIDE 20 — The result (3 sub-zones) ===
    s = new_slide(20, "PART D  ·  THREE SUB-ZONES, CONSISTENT")
    add_title(s, "Sub-zone 0 → 1 → 2. Quality improves monotonically.")
    add_subtitle(s, "Pooled K-Means centroids on Zone C (3,571 net samples, 6 features).")

    # Table of centroids
    rows, cols = 4, 6
    tbl_y = Inches(2.4)
    tbl_h = Inches(2.2)
    tbl = s.shapes.add_table(rows, cols, MARGIN_L, tbl_y, CONTENT_W, tbl_h).table
    widths = [1.3, 1.5, 1.6, 1.6, 1.4, 4.7]
    total = sum(widths)
    for i, w in enumerate(widths):
        tbl.columns[i].width = Inches(w * 12.13 / total)
    headers = ["Sub-zone", "Avg vsh", "Avg phit", "Avg perm (mD)", "Avg sw", "Tier"]
    for c, h in enumerate(headers):
        style_table_cell(tbl.cell(0, c), text=h, font_size=11, bold=True,
                         color=NAVY, bg=TABLE_HEADER_BG)

    centroids = [
        ("0", "0.510", "0.138", "~71",  "0.62", "Poor — high shale, low porosity, low perm", DANGER),
        ("1", "0.331", "0.201", "~231", "0.47", "Moderate — typical sand", ORANGE),
        ("2", "0.216", "0.246", "~744", "0.32", "Best — clean, high-perm drilling target", SUCCESS),
    ]
    for r, (sz, vsh, phit, perm, sw, tier, color) in enumerate(centroids, start=1):
        bg = TABLE_ROW_BG
        style_table_cell(tbl.cell(r, 0), text=sz, font_size=14, bold=True,
                         color=color, bg=bg, align=PP_ALIGN.CENTER)
        style_table_cell(tbl.cell(r, 1), text=vsh, font_size=11, color=INK,
                         bg=bg, align=PP_ALIGN.RIGHT)
        style_table_cell(tbl.cell(r, 2), text=phit, font_size=11, color=INK,
                         bg=bg, align=PP_ALIGN.RIGHT)
        style_table_cell(tbl.cell(r, 3), text=perm, font_size=11, color=INK,
                         bg=bg, align=PP_ALIGN.RIGHT)
        style_table_cell(tbl.cell(r, 4), text=sw, font_size=11, color=INK,
                         bg=bg, align=PP_ALIGN.RIGHT)
        style_table_cell(tbl.cell(r, 5), text=tier, font_size=11, color=color, bg=bg)

    # Bottom highlight
    add_quote_card(s, MARGIN_L, Inches(5.2), CONTENT_W, Inches(1.4), accent=GOLD)
    add_text(s, MARGIN_L + Inches(0.3), Inches(5.35), CONTENT_W - Inches(0.6),
             Inches(0.45),
             "All three sub-zones appear in all seven wells.",
             font_size=14, color=NAVY, bold=True)
    add_text(s, MARGIN_L + Inches(0.3), Inches(5.8), CONTENT_W - Inches(0.6),
             Inches(0.8),
             "21 (well × sub-zone) rows  ·  no missing combinations  ·  perm gap sub-zone 0 → 2 is 5-7× in every well.",
             font_size=11, color=INK, spacing=1.4, italic=True)

    # === SLIDE 21 — LOWO ARI ===
    s = new_slide(21, "PART D  ·  CROSS-WELL CONSISTENCY")
    add_title(s, "Validation by leave-one-well-out.")
    add_subtitle(s, "Fit on six wells, predict on the seventh, compare to pooled clustering.")

    y0 = Inches(2.4)
    # Big LOWO number
    add_text(s, MARGIN_L, y0, Inches(4.5), Inches(1.5),
             "0.991", font_size=110, color=NAVY, bold=True, spacing=1.0)
    add_text(s, MARGIN_L, y0 + Inches(1.55), Inches(4.5), Inches(0.4),
             "mean LOWO ARI", font_size=14, color=INK, bold=True)
    add_text(s, MARGIN_L, y0 + Inches(2.0), Inches(4.5), Inches(0.4),
             "Adjusted Rand Index", font_size=10, color=MUTED, italic=True)
    add_text(s, MARGIN_L, y0 + Inches(2.5), Inches(4.5), Inches(0.6),
             "1.0 = identical clustering. The structure is essentially invariant to which well is withheld.",
             font_size=10, color=MUTED, spacing=1.4)

    # Per-well grid
    rx = Inches(6.6)
    add_text(s, rx, y0, Inches(6.0), Inches(0.4),
             "Per-fold scores", font_size=10, color=GOLD, bold=True)
    lowo_data = [
        (1, "0.963"), (2, "0.992"), (3, "1.000"),
        (4, "0.984"), (5, "1.000"), (6, "1.000"),
        (7, "0.996"),
    ]
    cell_w = Inches(0.78)
    cell_h = Inches(0.78)
    cell_gap = Inches(0.08)
    cx = rx
    cy = y0 + Inches(0.5)
    for i, (well, val) in enumerate(lowo_data):
        color = SUCCESS if val == "1.000" else NAVY
        cb = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, cy, cell_w, cell_h)
        cb.fill.solid()
        cb.fill.fore_color.rgb = color
        cb.line.fill.background()
        add_text(s, cx, cy + Inches(0.08), cell_w, Inches(0.3),
                 f"w{well}", font_size=10, color=WHITE,
                 bold=True, align=PP_ALIGN.CENTER)
        add_text(s, cx, cy + Inches(0.4), cell_w, Inches(0.3),
                 val, font_size=11, color=WHITE,
                 bold=True, align=PP_ALIGN.CENTER)
        cx += cell_w + cell_gap

    # Insight box below the grid
    add_text(s, rx, y0 + Inches(1.6), Inches(6.0), Inches(0.4),
             "Three folds give exact 1.000.",
             font_size=12, color=NAVY, bold=True)
    add_text(s, rx, y0 + Inches(2.0), Inches(6.0), Inches(0.9),
             "Withholding wells 3, 5, or 6 produces the identical pooled clustering. The other four folds match within 0.04 of 1.0.",
             font_size=11, color=MUTED, spacing=1.4, italic=True)

    # Bottom callout
    add_quote_card(s, MARGIN_L, Inches(6.0), CONTENT_W, Inches(0.85), accent=NAVY)
    add_text(s, MARGIN_L + Inches(0.3), Inches(6.1), CONTENT_W - Inches(0.6),
             Inches(0.7),
             "Operational meaning: sub-zone 0 in well 1 has the same characteristic vsh, phit, perm, and sw as sub-zone 0 in well 7. Same rock, different depths.",
             font_size=11, color=INK, italic=True, spacing=1.4)

    # === SLIDE 22 — Drilling target ===
    s = new_slide(22, "PART D  ·  THE DRILLING IMPLICATION")
    add_title(s, "Sub-zone 2 — the natural drilling target.")
    add_subtitle(s, "The clustering isn't academic. It changes the per-foot economics.")

    y0 = Inches(2.3)
    # Hero stat
    add_text(s, MARGIN_L, y0, CONTENT_W, Inches(1.5),
             "51%  of  kh",
             font_size=96, color=NAVY, bold=True, align=PP_ALIGN.CENTER, spacing=1.0)
    add_text(s, MARGIN_L, y0 + Inches(1.6), CONTENT_W, Inches(0.5),
             "in just 29% of Zone C's thickness.",
             font_size=20, color=INK, italic=True, align=PP_ALIGN.CENTER)

    # 4 supporting stats
    y1 = Inches(5.0)
    stats = [
        ("315.6 m", "field-wide thickness"),
        ("~1,100 mD", "avg permeability"),
        ("~7×",     "perm gap vs sub-zone 0"),
        ("7 / 7",   "wells where present"),
    ]
    for i, (val, lbl) in enumerate(stats):
        x = MARGIN_L + Inches(3.05) * i
        add_text(s, x, y1, Inches(3.0), Inches(0.6),
                 val, font_size=24, color=NAVY, bold=True)
        add_text(s, x, y1 + Inches(0.65), Inches(3.0), Inches(0.3),
                 lbl, font_size=10, color=MUTED)

    add_text(s, MARGIN_L, Inches(6.8), CONTENT_W, Inches(0.4),
             "→  Targeting sub-zone 2 captures roughly twice the flow per metre drilled vs an undifferentiated Zone C target.",
             font_size=11, color=NAVY, italic=True)

    # === SLIDE 23 — Zone B negative result ===
    s = new_slide(23, "PART D  ·  THE NEGATIVE RESULT")
    add_title(s, "We ran Zone B as a controlled experiment.")
    add_subtitle(s, "It failed in a specific, predictable way — and that failure is itself a finding.")

    y0 = Inches(2.4)
    # The smoking gun
    add_text(s, MARGIN_L, y0, CONTENT_W, Inches(0.4),
             "Zone B sub-zone centroids — log_perm values:",
             font_size=12, color=MUTED, bold=True)

    # Two centroid values side by side
    box_w = Inches(5.9)
    box_h = Inches(2.2)
    box1 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN_L, Inches(2.95), box_w, box_h)
    box1.fill.solid()
    box1.fill.fore_color.rgb = WHITE
    box1.line.color.rgb = TABLE_BORDER
    box1.line.width = Pt(0.5)
    add_text(s, MARGIN_L + Inches(0.3), Inches(3.1), box_w - Inches(0.6),
             Inches(0.3), "SUB-ZONE 1", font_size=10, color=MUTED, bold=True)
    add_text(s, MARGIN_L + Inches(0.3), Inches(3.4), box_w - Inches(0.6),
             Inches(1.0), "4.175938", font_size=48, color=NAVY, bold=True, spacing=1.0)
    add_text(s, MARGIN_L + Inches(0.3), Inches(4.45), box_w - Inches(0.6),
             Inches(0.3), "≈ 14,985 mD", font_size=11, color=MUTED, italic=True)

    box2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.0), Inches(2.95), box_w, box_h)
    box2.fill.solid()
    box2.fill.fore_color.rgb = WHITE
    box2.line.color.rgb = TABLE_BORDER
    box2.line.width = Pt(0.5)
    add_text(s, Inches(7.3), Inches(3.1), box_w - Inches(0.6), Inches(0.3),
             "SUB-ZONE 2", font_size=10, color=MUTED, bold=True)
    add_text(s, Inches(7.3), Inches(3.4), box_w - Inches(0.6), Inches(1.0),
             "4.176091", font_size=48, color=NAVY, bold=True, spacing=1.0)
    add_text(s, Inches(7.3), Inches(4.45), box_w - Inches(0.6), Inches(0.3),
             "≈ 14,995 mD", font_size=11, color=MUTED, italic=True)

    # Gap callout
    add_quote_card(s, MARGIN_L, Inches(5.4), CONTENT_W, Inches(0.95),
                   bg=HIGHLIGHT_BG, accent=WARNING)
    add_text(s, MARGIN_L + Inches(0.3), Inches(5.5), Inches(4.5), Inches(0.4),
             "Difference: 0.00015",
             font_size=16, color=WARNING, bold=True)
    add_text(s, MARGIN_L + Inches(0.3), Inches(5.9), CONTENT_W - Inches(0.6), Inches(0.5),
             "10 mD on a permeability spanning 5 orders of magnitude. The tool can't tell them apart — and so neither can clustering.",
             font_size=11, color=INK, spacing=1.4, italic=True)

    add_text(s, MARGIN_L, Inches(6.65), CONTENT_W, Inches(0.4),
             "→  Saturation defeats clustering on a saturated feature. Confirms what Lorenz already suggested.",
             font_size=11, color=NAVY, italic=True)

    # === SLIDE 24 — Lessons (key takeaways) ===
    s = new_slide(24, "RECAP  ·  THE LESSONS")
    add_title(s, "Three things we took away.")
    add_subtitle(s, "Methodology lessons that apply beyond this dataset.")

    y0 = Inches(2.4)
    lessons = [
        ("01", "Discovery first, computation second",
         "Saturation could only be found by interrogating the data, not by accepting it. Every downstream finding rests on this single discovery."),
        ("02", "Bound conclusions when calibration is missing",
         "Without core or production data, sensitivity sweeps and uncertainty bounds replace single point estimates. Robust findings survive any reasonable cutoff."),
        ("03", "Negative results matter as much as positive ones",
         "Zone B's clustering failure is as informative as Zone C's success. It confirms — at a different level — what Lorenz already suggested."),
    ]
    for i, (num, head, body) in enumerate(lessons):
        ry = y0 + Inches(1.35) * i
        add_text(s, MARGIN_L, ry, Inches(0.7), Inches(0.7),
                 num, font_size=32, color=GOLD, bold=True)
        add_text(s, MARGIN_L + Inches(0.85), ry, Inches(11), Inches(0.4),
                 head, font_size=15, color=NAVY, bold=True)
        add_text(s, MARGIN_L + Inches(0.85), ry + Inches(0.45),
                 Inches(11.3), Inches(0.85), body, font_size=11, color=MUTED, spacing=1.4)

    # === SLIDE 25 — What's next ===
    s = new_slide(25, "RECAP  ·  WHAT'S NEXT")
    add_title(s, "What another two weeks would buy.")
    add_subtitle(s, "Three concrete next steps to lift the conclusions from internal-consistency to external validation.")

    y0 = Inches(2.4)
    nexts = [
        ("Core data validation",
         "Compare sub-zone boundaries to described lithology in core photos and measured core perm. Gold-standard external validation.",
         "Highest value, requires core lab access"),
        ("Production log overlay",
         "Match PLT/spinner inflow profiles against sub-zone assignments. Tests whether sub-zone 2 actually produces more per metre.",
         "Operational validation — confirms drilling implication"),
        ("Saturation-corrected kh bounds",
         "Estimate the true Zone B kh range using sandstone perm priors. Pair the current lower bound with an order-of-magnitude upper bound.",
         "Resolves the largest single uncertainty in the deliverable"),
    ]
    for i, (head, body, footer) in enumerate(nexts):
        ry = y0 + Inches(1.35) * i
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, MARGIN_L, ry + Inches(0.18),
                                  Inches(0.18), Inches(0.18))
        dot.fill.solid()
        dot.fill.fore_color.rgb = GOLD
        dot.line.fill.background()
        add_text(s, MARGIN_L + Inches(0.5), ry, Inches(11), Inches(0.4),
                 head, font_size=14, color=NAVY, bold=True)
        add_text(s, MARGIN_L + Inches(0.5), ry + Inches(0.4), Inches(11.5), Inches(0.5),
                 body, font_size=11, color=INK, spacing=1.4)
        add_text(s, MARGIN_L + Inches(0.5), ry + Inches(0.95), Inches(11.5), Inches(0.3),
                 footer, font_size=10, color=MUTED, italic=True)

    # === SLIDE 26 — Thank you ===
    s = prs.slides.add_slide(blank)
    add_background(s)
    add_left_accent_strip(s, NAVY)
    add_text(s, Inches(0.85), Inches(2.5), Inches(11), Inches(0.4),
             "RESERVOIR ANALYTICS  ·  MAY 2026", font_size=11,
             color=GOLD, bold=True)
    add_text(s, Inches(0.85), Inches(3.0), Inches(11), Inches(1.5),
             "Thank you.", font_size=72, color=INK, bold=True, spacing=1.0)
    add_horizontal_rule(s, Inches(4.7), x=Inches(0.85), w=Inches(2.5),
                        color=GOLD, height=Pt(1.5))
    add_text(s, Inches(0.85), Inches(4.95), Inches(11), Inches(0.4),
             "Kamil Muradli", font_size=18, color=INK, bold=True)
    add_text(s, Inches(0.85), Inches(5.4), Inches(11), Inches(0.4),
             "eiGroup Associate Data Scientist Assessment", font_size=12, color=MUTED)
    add_text(s, Inches(0.85), Inches(5.85), Inches(11), Inches(0.4),
             "Submitted May 2026", font_size=10, color=PAGE_NUM)
    add_text(s, Inches(0.85), Inches(6.4), Inches(11), Inches(0.4),
             "Questions and detailed walkthrough documents available in /outputs/reports/",
             font_size=11, color=NAVY, italic=True)
    add_page_number(s, 26)

    # Save
    prs.save(str(OUT_PATH))
    print(f"Wrote {OUT_PATH}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build_presentation()